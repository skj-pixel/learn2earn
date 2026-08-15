# -*- coding: utf-8 -*-
# 🔍 [语法] UTF-8 编码
# 🔍 [语法] docstring
# 🔍 [作用] 文件用途——PPT 生成工具
"""生成学赚7步法PPT - 初赛/复赛/决赛"""

# 🔍 [语法] 标准库
import os
# 🔍 [语法] python-pptx
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.oxml.ns import qn

# ---- 颜色常量 ----
# 🔍 [语法] 模块级常量
# 🔍 [作用] 复用 common.py 的颜色调色板
# 🔍 [陷阱] ⚠️ 与 common.py 重复定义（应 import common 替代）
LIGHT_BG = RGBColor(0xFF, 0xFF, 0xFF); INDIGO = RGBColor(0x63, 0x66, 0xF1); PURPLE = RGBColor(0x8B, 0x5C, 0xF6)
LIGHT_INDIGO = RGBColor(0xE0, 0xE7, 0xFE); LIGHT_PURPLE = RGBColor(0xED, 0xE9, 0xFE)
DARK_TEXT = RGBColor(0x1F, 0x29, 0x37); GRAY_TEXT = RGBColor(0x6B, 0x72, 0x80)
ACCENT_CYAN = RGBColor(0x22, 0xD3, 0xEE); SUCCESS_GREEN = RGBColor(0x10, 0xB9, 0x81)
WARN_ORANGE = RGBColor(0xF5, 0x9E, 0x0B); DANGER_RED = RGBColor(0xEF, 0x44, 0x44); WHITE = RGBColor(0xFF, 0xFF, 0xFF)
LG = RGBColor(0xD1, 0xFA, 0xE5); LC = RGBColor(0xEC, 0xFE, 0xFF); LY = RGBColor(0xFE, 0xF3, 0xC7); LR = RGBColor(0xFE, 0xE2, 0xE2)
DARK_BG = RGBColor(0x0F, 0x17, 0x2A); DARK_CARD = RGBColor(0x1E, 0x29, 0x3B); DARK_CARD2 = RGBColor(0x33, 0x41, 0x55)
BI = RGBColor(0x81, 0x8C, 0xF8); BC = RGBColor(0x22, 0xD3, 0xEE); BP = RGBColor(0xA7, 0x8B, 0xFA)
WT = RGBColor(0xF8, 0xFA, 0xFC); GL = RGBColor(0x94, 0xA3, 0xB8)

# 🔍 [语法] 字体 + 单位缩写
# 🔍 [作用] 全局常量（与 common.py 一致）
FONT = "微软雅黑"
I = Inches
P = Pt


# 🔍 [语法] 简写工具函数
# 🔍 [作用] 与 common.py 类似（重复定义以独立使用）
def set_bg(s, c): f = s.background.fill; f.solid(); f.fore_color.rgb = c

def _run(r, sz=14, b=False, c=DARK_TEXT):
    r.font.name = FONT; r.font.size = P(sz); r.font.bold = b; r.font.color.rgb = c
    rp = r._r.get_or_add_rPr()
    for e in rp.findall(qn('a:ea')): rp.remove(e)
    rp.append(rp.makeelement(qn('a:ea'), {'typeface': FONT}))


def tx(s, l, t, w, h, txt, sz=14, b=False, c=DARK_TEXT, *, b2=False, a=PP_ALIGN.LEFT, ac=MSO_ANCHOR.TOP, ls=1.3):
    """添加文本框（支持中文换行 \\n）"""
    tb = s.shapes.add_textbox(l, t, w, h); tf = tb.text_frame; tf.word_wrap = True
    tf.margin_left = P(4); tf.margin_right = P(4); tf.margin_top = P(2); tf.margin_bottom = P(2); tf.vertical_anchor = ac
    lines = txt.split('\n') if isinstance(txt, str) else txt
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = a; p.line_spacing = ls
        r = p.add_run(); r.text = line; _run(r, sz, b, c)
    return tb


def bul(s, l, t, w, h, items, sz=14, c=DARK_TEXT, bc=None, ls=1.4):
    """添加项目符号列表"""
    tb = s.shapes.add_textbox(l, t, w, h); tf = tb.text_frame; tf.word_wrap = True
    tf.margin_left = P(4); tf.margin_right = P(4); bc = bc or c
    for i, it in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = PP_ALIGN.LEFT; p.line_spacing = ls; p.space_after = P(4)
        r0 = p.add_run(); r0.text = "● "; _run(r0, sz, True, bc)
        if isinstance(it, tuple):
            hh, bd = it; r1 = p.add_run(); r1.text = hh; _run(r1, sz, True, c)
            r2 = p.add_run(); r2.text = bd; _run(r2, sz, False, c)
        else:
            r = p.add_run(); r.text = it; _run(r, sz, False, c)
    return tb


def shp(s, l, t, w, h, f, ln=None, sh=MSO_SHAPE.RECTANGLE):
    """添加形状"""
    sh = s.shapes.add_shape(sh, l, t, w, h); sh.fill.solid(); sh.fill.fore_color.rgb = f
    if ln is None: sh.line.fill.background()
    else: sh.line.color.rgb = ln; sh.line.width = P(1)
    sp = sh.fill._xPr
    if sp.find(qn('a:effectLst')) is None: sp.append(sp.makeelement(qn('a:effectLst'), {}))
    return sh


def R(s, l, t, w, h, f, ln=None): return shp(s, l, t, w, h, f, ln, MSO_SHAPE.RECTANGLE)
def RR(s, l, t, w, h, f, ln=None): return shp(s, l, t, w, h, f, ln, MSO_SHAPE.ROUNDED_RECTANGLE)
def OV(s, l, t, w, h, f): return shp(s, l, t, w, h, f, None, MSO_SHAPE.OVAL)
def AR(s, l, t, w, h, c): return shp(s, l, t, w, h, c, None, MSO_SHAPE.RIGHT_ARROW)


def lt(s, title, sub=None):
    """浅色页头"""
    R(s, I(0), I(0), I(13.333), I(0.18), INDIGO)
    tx(s, I(0.5), I(0.3), I(12.3), I(0.6), title, 28, DARK_TEXT, b=True, ac=MSO_ANCHOR.MIDDLE)
    if sub: tx(s, I(0.5), I(0.9), I(12.3), I(0.4), sub, 14, GRAY_TEXT)
    R(s, I(0), I(7.35), I(13.333), I(0.05), LIGHT_INDIGO)


def dt(s, title, sub=None):
    """深色页头"""
    R(s, I(0), I(0), I(13.333), I(0.12), BI)
    tx(s, I(0.5), I(0.25), I(12.3), I(0.6), title, 26, WT, b=True, ac=MSO_ANCHOR.MIDDLE)
    if sub: tx(s, I(0.5), I(0.82), I(12.3), I(0.35), sub, 13, GL)
    R(s, I(0), I(7.35), I(13.333), I(0.04), DARK_CARD2)


def pn(s, n, T, dark=False):
    """页脚（n / T）"""
    tx(s, I(11.8), I(7.0), I(1.3), I(0.3), f"{n} / {T}", 10, c=GL if dark else GRAY_TEXT, a=PP_ALIGN.RIGHT)


def tbl(s, l, t, w, h, data, cw=None, hc=INDIGO, ht=WHITE, bt=DARK_TEXT, alt=LIGHT_INDIGO, dark=False, fs=12):
    """添加表格"""
    rs = len(data); cs = len(data[0])
    t = s.shapes.add_table(rs, cs, l, t, w, h).table
    if cw:
        S = sum(cw)
        for i, c in enumerate(cw): t.columns[i].width = int(w * c / S)
    for ri, row in enumerate(data):
        for ci, v in enumerate(row):
            ce = t.cell(ri, ci)
            ce.margin_left = P(6); ce.margin_right = P(6); ce.margin_top = P(4); ce.margin_bottom = P(4)
            ce.vertical_anchor = MSO_ANCHOR.MIDDLE
            if ri == 0:
                ce.fill.solid(); ce.fill.fore_color.rgb = hc; tc = ht; bo = True
            else:
                ce.fill.solid()
                ce.fill.fore_color.rgb = (DARK_CARD if ri % 2 == 1 else DARK_CARD2) if dark else (alt if ri % 2 == 0 else LIGHT_BG)
                tc = bt; bo = False
            tf = ce.text_frame; tf.word_wrap = True; tf.clear()
            p = tf.paragraphs[0]
            p.alignment = PP_ALIGN.CENTER if (ci == 0 or ri == 0) else PP_ALIGN.LEFT
            p.line_spacing = 1.2
            r = p.add_run(); r.text = str(v); _run(r, fs, bo, tc)


def new():
    """新建 16:9 PPT"""
    p = Presentation(); p.slide_width = I(13.333); p.slide_height = I(7.5); return p


def B(p): return p.slides.add_slide(p.slide_layouts[6])
